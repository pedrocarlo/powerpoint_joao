from tkinter.filedialog import askdirectory
from typing import List, Dict
import os
import pptx.shapes.picture
from pptx import Presentation
from pptx.util import Inches, Pt

BLANK_SLIDE = 6
slide_height = Inches(7.5).emu  # Default Slide Height
slide_width = Inches(13.34).emu  # Default Slide Width


class Slide:
    def __init__(self, template_file: str, images: List[str] = None, texts: List[str] = None):
        if images is not None:
            self.images = images
        else:
            self.images = None
        self.texts = texts  # stores the texts for the textboxes
        self.template_file = template_file
        self.template = Presentation(template_file)
        # read template and determine number of pictures and textboxes
        self.num_images, self.num_text = self.number_pictures_text()

    def update_images(self, images: List[str]):
        self.images = images

    def number_pictures_text(self):
        img_count = 0
        text_count = 0
        for slide in self.template.slides:
            for shape in slide.shapes:
                if isinstance(shape,
                              pptx.shapes.autoshape.Shape) and shape.has_text_frame:  # In this case the autoshape in question is a textbox
                    text_count += 1
                if isinstance(shape, pptx.shapes.picture.Picture):  # If object is a Picture
                    img_count += 1
        return img_count, text_count

    def add_slide(self, presentation: Presentation):
        add_pictures_text(presentation, Presentation(self.template_file), self)

    def get_images(self):
        return self.images


def new_presentation(title):
    new_prs = Presentation()
    new_prs.slide_height = slide_height
    new_prs.slide_width = slide_width
    return new_prs, title


def add_pictures_text(presentation: Presentation, template: Presentation, new_slide: Slide):
    blank_slide_layout = presentation.slide_layouts[6]
    new_pres_slide = presentation.slides.add_slide(blank_slide_layout)
    img_count = 0
    text_count = 0
    for slide in template.slides:
        shapes = list(slide.shapes)
        shapes.reverse()
        for shape in shapes:
            if isinstance(shape,
                          pptx.shapes.autoshape.Shape) and shape.has_text_frame:  # In this case the autoshape in
                # question is a textbox
                try:
                    left = shape.left
                    top = shape.top
                    height = shape.height
                    width = shape.width
                    txBox = new_pres_slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    text_list = new_slide.texts[text_count].split("\n")
                    for i in range(len(text_list)):
                        tf.add_paragraph()
                    for i in range(len(tf.paragraphs)):
                        p = tf.paragraphs[i]
                        run = p.add_run()
                        p.alignment = shape.text_frame.paragraphs[i].alignment
                        font = run.font
                        font.name = "Franklin Gothic Demi"
                        font.size = Pt(24)
                        font.bold = True
                        font.italic = False
                        run.text = text_list[i]

                    tf.text = new_slide.texts[text_count]  # TODO CHANGE TEXT HERE TO BE FROM SLIDE

                except IndexError:  # If user does not supply enough text boxes the code will still run
                    continue
                text_count += 1
            if isinstance(shape, pptx.shapes.picture.Picture):  # If object is a Picture
                try:
                    left = shape.left
                    top = shape.top
                    height = shape.height
                    new_pres_slide.shapes.add_picture(new_slide.images[img_count],
                                                      left, top, height=height)
                except (IndexError, FileNotFoundError,
                        FileExistsError):  # if user does not submit enough photos it will still run
                    continue
                img_count += 1


def save_presentation(presentation: Presentation, slides, title, template_dir):
    for slide in slides.values():
        add_pictures_text(presentation, slide.template, slide)
    desktop = os.path.join(os.path.join(os.path.expanduser('~')), 'Desktop')  # Desktop
    presentation_dir = askdirectory(title='Select Template Directory', initialdir=os.path.dirname(template_dir))
    full_pres_dir_path = presentation_dir
    # if not (os.path.exists(full_pres_dir_path) and os.path.isdir(full_pres_dir_path)):
    #     os.mkdir(full_pres_dir_path)  # create directory if it does not exist
    presentation.save(full_pres_dir_path + "/" + title + ".pptx")
